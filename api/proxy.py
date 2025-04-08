import os
import requests
from flask import (
    Flask,
    request,
    Response,
    abort,
    jsonify,
    render_template_string,
    send_file,
    redirect,
    g,  # Added g
)
from requests_ntlm import HttpNtlmAuth
from io import BytesIO
import concurrent.futures
import PyPDF2
import hashlib
import os
import time
from werkzeug.http import parse_range_header
import threading
from functools import lru_cache
import gzip
import redis
from dotenv import load_dotenv
import base64
import logging  # Added
import json  # Added
from time import perf_counter  # Added
from datetime import datetime, timezone  # Added timezone
import traceback  # Added
import sys  # Added
import atexit  # Added

# Load environment variables
load_dotenv()

# --- Append parent directory if needed ---
# (Keep if necessary for your project structure)
# sys.path.append(os.path.join(os.path.dirname(__file__), ".."))

app = Flask(__name__)

# Set up logging using the built-in logging module.
# Use a more specific logger name if desired
logging.basicConfig(
    level=logging.INFO,  # Changed default level to INFO
    format="%(asctime)s - %(levelname)s - [%(threadName)s] - %(message)s",
)
logger = logging.getLogger("proxy_extractor_api")  # Use a specific logger

# Configuration - REDIS OPTIMIZED VERSION
REDIS_URL = os.environ.get("REDIS_URL", "redis://localhost:6379")
CACHE_EXPIRY = 1800  # 30 minutes in seconds
MAX_CHUNK_SIZE = 1024 * 1024  # 1MB chunk size for Redis storage
CHUNK_SIZE = 262144  # 256KB chunks for better streaming performance
ENABLE_BACKGROUND_THREADS = False  # Set to True on dedicated servers, False on Vercel

# Initialize Redis client
try:
    # Use decode_responses=True for easier string handling in metadata
    redis_client = redis.from_url(REDIS_URL, decode_responses=True)
    redis_client.ping()  # Test connection
    logger.info(f"Connected to Redis at {REDIS_URL}")
except Exception as e:
    logger.critical(f"Failed to connect to Redis: {e}", exc_info=True)
    redis_client = None  # Set to None to indicate failure

# --- Logging Constants ---
API_LOG_KEY = "api_logs"  # Specific Redis key for proxy/extractor logs
MAX_LOG_ENTRIES = 5000

# --- Thread Pool for Background Logging ---
log_executor = concurrent.futures.ThreadPoolExecutor(
    max_workers=5, thread_name_prefix="LogThread"
)


# --- Background Logging Task ---
def _log_to_redis_task(log_entry_dict):
    """Internal task to write logs to Redis asynchronously."""
    if not redis_client:
        print(
            f"[{threading.current_thread().name}] Log Error: Redis client not available. Skipping log.",
            file=sys.stderr,
        )
        return
    try:
        log_entry_json = json.dumps(log_entry_dict)
        # Use pipeline for atomic push+trim
        pipe = redis_client.pipeline()
        pipe.lpush(API_LOG_KEY, log_entry_json)
        pipe.ltrim(API_LOG_KEY, 0, MAX_LOG_ENTRIES - 1)
        pipe.execute()
    except redis.exceptions.ConnectionError as e:
        print(
            f"[{threading.current_thread().name}] Log Error: Redis connection error: {e}",
            file=sys.stderr,
        )
    except TypeError as e:
        print(
            f"[{threading.current_thread().name}] Log Error: Failed to serialize log entry to JSON: {e}",
            file=sys.stderr,
        )
        problematic_items = {
            k: repr(v)
            for k, v in log_entry_dict.items()
            if not isinstance(v, (str, int, float, bool, list, dict, type(None)))
        }
        print(
            f"[{threading.current_thread().name}] Log Error: Problematic items: {problematic_items}",
            file=sys.stderr,
        )
    except Exception as e:
        print(
            f"[{threading.current_thread().name}] Log Error: Failed to write log to Redis: {e}",
            file=sys.stderr,
        )
        print(
            f"[{threading.current_thread().name}] Log Error: Failed entry (partial): user={log_entry_dict.get('username')}, endpoint={log_entry_dict.get('endpoint')}, status={log_entry_dict.get('status_code')}",
            file=sys.stderr,
        )


# Connection pool for reusing connections
session = requests.Session()
adapter = requests.adapters.HTTPAdapter(
    pool_connections=10, pool_maxsize=20, max_retries=3, pool_block=False
)
session.mount("http://", adapter)
session.mount("https://", adapter)

# Modified prefetch implementation for serverless environment
prefetch_queue = []
prefetch_lock = threading.Lock()


def get_cache_key(url, username):
    """Generate a unique cache key for a URL and username"""
    key = f"proxy:{url}:{username}"
    # Use hexdigest for shorter keys if needed, but full key is more readable
    # return hashlib.md5(key.encode()).hexdigest()
    return key  # Using readable key


def get_metadata_key(cache_key):
    """Get the key for file metadata in Redis"""
    return f"metadata:{cache_key}"


def is_cached(cache_key):
    """Check if file is cached and valid in Redis"""
    if not redis_client:
        return False
    try:
        metadata_key = get_metadata_key(cache_key)
        # Check existence first for efficiency
        if not redis_client.exists(metadata_key):
            return False
        metadata = redis_client.hgetall(metadata_key)
        if not metadata:  # Should not happen if exists check passed, but defensive
            return False
        # decode_responses=True means keys/values are already strings
        expiry_time = float(metadata.get("expiry", 0))
        is_valid = time.time() < expiry_time
        if not is_valid:
            logger.info(f"Cache expired for key: {cache_key}")
            # Optionally delete expired keys here or rely on Redis TTL
            # redis_client.delete(metadata_key) # Be careful with concurrent access
        return is_valid
    except redis.exceptions.ConnectionError as e:
        logger.error(f"Redis connection error checking cache for {cache_key}: {e}")
        return False
    except Exception as e:
        logger.error(f"Error checking cache for {cache_key}: {e}", exc_info=True)
        return False


def save_to_cache(cache_key, content):
    """Save content to Redis cache with chunking for large files"""
    if not redis_client:
        logger.warning("Redis not available, skipping cache save")
        return
    try:
        metadata = {
            "size": len(content),
            "chunks": (len(content) + MAX_CHUNK_SIZE - 1) // MAX_CHUNK_SIZE,
            "expiry": time.time() + CACHE_EXPIRY,
            "created": time.time(),
        }
        metadata_key = get_metadata_key(cache_key)

        # Store metadata using hmset (or hset in newer redis-py versions)
        redis_client.hmset(metadata_key, metadata)
        redis_client.expire(metadata_key, CACHE_EXPIRY)

        # Store content in chunks
        num_chunks = metadata["chunks"]
        pipe = redis_client.pipeline() # Use pipeline for efficiency
        for i in range(num_chunks):
            chunk = content[i * MAX_CHUNK_SIZE : (i + 1) * MAX_CHUNK_SIZE]
            chunk_key = f"chunk:{cache_key}:{i}"
            # Use base64 encoding for binary data
            encoded_chunk = base64.b64encode(chunk).decode(
                "ascii"
            )  # Store as string
            pipe.setex(chunk_key, CACHE_EXPIRY, encoded_chunk)
        pipe.execute() # Execute all chunk saves

        logger.info(
            f"Saved to Redis cache: {cache_key} ({len(content)} bytes in {num_chunks} chunks)"
        )
    except redis.exceptions.ConnectionError as e:
        logger.error(f"Redis connection error saving cache for {cache_key}: {e}")
    except Exception as e:
        logger.error(f"Redis cache save error for {cache_key}: {e}", exc_info=True)


def get_from_cache(cache_key):
    """Get content from Redis cache, assembling chunks if needed"""
    if not redis_client:
        return None
    try:
        metadata_key = get_metadata_key(cache_key)
        metadata = redis_client.hgetall(metadata_key)
        if not metadata:
            return None

        # decode_responses=True means keys/values are already strings
        num_chunks = int(metadata.get("chunks", 0))
        if num_chunks == 0:
            logger.warning(f"Metadata for {cache_key} indicates 0 chunks.")
            return b"" # Return empty bytes if size was 0

        content = bytearray()
        pipe = redis_client.pipeline() # Use pipeline for efficient fetching
        chunk_keys = [f"chunk:{cache_key}:{i}" for i in range(num_chunks)]
        for key in chunk_keys:
            pipe.get(key)
        encoded_chunks = pipe.execute() # Fetch all chunks

        for i, encoded_chunk in enumerate(encoded_chunks):
            if not encoded_chunk:
                logger.error(f"Missing chunk {i} for key {cache_key}")
                # Optionally delete metadata if cache is inconsistent
                # redis_client.delete(metadata_key)
                return None
            try:
                # Decode from base64 (encoded_chunk is already string)
                chunk = base64.b64decode(encoded_chunk)
                content.extend(chunk)
            except Exception as decode_err:
                 logger.error(f"Error decoding base64 chunk {i} for {cache_key}: {decode_err}")
                 return None

        # Verify size if needed
        expected_size = int(metadata.get("size", -1))
        if expected_size != -1 and len(content) != expected_size:
             logger.error(f"Cache size mismatch for {cache_key}. Expected {expected_size}, got {len(content)}. Cache might be corrupt.")
             # Optionally delete inconsistent cache
             # redis_client.delete(metadata_key, *chunk_keys)
             return None

        logger.info(f"Retrieved from Redis cache: {cache_key} ({len(content)} bytes)")
        return bytes(content)
    except redis.exceptions.ConnectionError as e:
        logger.error(f"Redis connection error reading cache for {cache_key}: {e}")
        return None
    except Exception as e:
        logger.error(f"Redis cache read error for {cache_key}: {e}", exc_info=True)
        return None


def prefetch_worker():
    """Background worker to prefetch files"""
    while True:
        item = None
        try:
            with prefetch_lock:
                if prefetch_queue:
                    item = prefetch_queue.pop(0)

            if item:
                username, password, file_url = item
                cache_key = get_cache_key(file_url, username)

                if not is_cached(cache_key):
                    logger.info(f"Prefetching: {file_url}")
                    headers = {
                        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
                        "Accept": "*/*",
                        "Accept-Encoding": "gzip, deflate",
                    }
                    try:
                        response = session.get(
                            file_url,
                            auth=HttpNtlmAuth(username, password),
                            headers=headers,
                            timeout=30, # Shorter timeout for prefetch
                        )
                        if response.status_code == 200:
                            save_to_cache(cache_key, response.content)
                        else:
                             logger.warning(f"Prefetch failed for {file_url}: Status {response.status_code}")
                    except Exception as fetch_err:
                         logger.error(f"Prefetch download error for {file_url}: {fetch_err}")
            else:
                # No items in queue, sleep longer
                time.sleep(5)

        except Exception as e:
            logger.error(f"Prefetch worker error: {e}", exc_info=True)
            # Avoid busy-looping on persistent errors
            time.sleep(10)

        # Small delay even if item was processed
        time.sleep(0.1)


# Only start prefetch thread if enabled and not on serverless
if ENABLE_BACKGROUND_THREADS and redis_client: # Also check if redis is available
    prefetch_thread = threading.Thread(target=prefetch_worker, daemon=True)
    prefetch_thread.start()
    logger.info("Prefetch thread started")
elif not redis_client:
     logger.warning("Redis not available, prefetch thread not started.")
else:
    logger.info("Background threads disabled, prefetch thread not started.")


# --- Request Hooks for Logging ---
@app.before_request
def before_request_func():
    """Initialize request context for timing and logging."""
    g.start_time = perf_counter()
    g.request_time = datetime.now(timezone.utc)
    g.username = None # Will be set in the view function if possible
    g.log_outcome = "unknown" # Default outcome
    g.log_error_message = None # Default error message

@app.after_request
def after_request_logger(response):
    """Gathers log info, handles User-Agent robustly, and submits the logging task asynchronously."""
    # Skip logging for the root/test form page or OPTIONS requests
    if request.method == 'OPTIONS' or request.path == '/':
        return response

    elapsed_ms = (perf_counter() - g.start_time) * 1000

    # --- Robust User-Agent Handling ---
    ua_string_from_parsed = None
    ua_parse_error = False
    raw_ua_header = request.headers.get('User-Agent')

    try:
        if request.user_agent:
            ua_string_from_parsed = request.user_agent.string
    except Exception as e:
        ua_parse_error = True
        logger.error(f"Proxy API UA: Error accessing request.user_agent.string: {e}", exc_info=True)

    final_user_agent = ua_string_from_parsed if ua_string_from_parsed else raw_ua_header if raw_ua_header else "Unknown"
    if ua_parse_error and not raw_ua_header:
        final_user_agent = "Unknown (Parsing Error)"
    # --- End User-Agent Handling ---

    # Prepare Log Entry
    username = getattr(g, 'username', None) # Get username if set in route
    outcome = getattr(g, 'log_outcome', 'unknown')
    error_message = getattr(g, 'log_error_message', None)

    # Mask password in request arguments (query parameters for these endpoints)
    request_args = request.args.to_dict()
    if 'password' in request_args:
        request_args['password'] = '********'

    ip_address = request.headers.get('X-Forwarded-For', request.remote_addr) or "Unknown"

    log_entry = {
        "endpoint": request.path,
        "error_message": error_message,
        "ip_address": ip_address,
        "method": request.method,
        "outcome": outcome,
        "request_args": request_args, # Log query args
        "request_timestamp_utc": g.request_time.isoformat(),
        "response_size_bytes": response.content_length,
        "response_timestamp_utc": datetime.now(timezone.utc).isoformat(),
        "status_code": response.status_code,
        "time_elapsed_ms": round(elapsed_ms, 2),
        "user_agent": final_user_agent,
        "username": username,
    }

    # Submit logging task
    try:
        log_executor.submit(_log_to_redis_task, log_entry)
    except Exception as e:
        logger.exception(f"CRITICAL: Failed to submit log task to executor: {e}")

    return response

@app.after_request
def add_cors_headers(response):
    """Add CORS headers."""
    # This runs *after* after_request_logger
    response.headers["Access-Control-Allow-Origin"] = "*"
    response.headers["Access-Control-Allow-Headers"] = "Content-Type, Authorization, Range" # Added Range
    response.headers["Access-Control-Allow-Methods"] = "GET, POST, OPTIONS" # Keep original methods
    # Expose headers needed for range requests
    response.headers["Access-Control-Expose-Headers"] = "Content-Length, Content-Range, Accept-Ranges"
    return response


# --- Routes ---

@app.route("/")
def index():
    # A simple HTML page for testing the endpoints (logging skipped via after_request)
    html = """
    <!doctype html>
    <html>
    <head>
      <title>Test Proxy/Extractor API</title>
       <style> body { font-family: sans-serif; } label { display: block; margin-top: 1em; } input { width: 90%; max-width: 500px; } </style>
    </head>
    <body>
      <h1>Test API Endpoints</h1>

      <h2>Proxy Endpoint</h2>
      <form id="proxyForm">
        <label>Username:</label>
        <input type="text" id="username_proxy" name="username" value="GUC\\your_username"><br>
        <label>Password:</label>
        <input type="password" id="password_proxy" name="password"><br>
        <label>File URL:</label>
        <input type="text" id="fileUrl_proxy" name="fileUrl" placeholder="e.g., https://cms.guc.edu.eg/apps/.../file.pdf"><br>
        <button type="button" onclick="testProxy()">Test Proxy Download</button>
      </form>
      <div id="proxyResult"></div>

      <h2>Extractor Endpoint</h2>
      <form id="extractForm">
        <label>Username:</label>
        <input type="text" id="username_extract" name="username" value="GUC\\your_username"><br>
        <label>Password:</label>
        <input type="password" id="password_extract" name="password"><br>
        <label>File URL:</label>
        <input type="text" id="fileUrl_extract" name="fileUrl" placeholder="e.g., https://cms.guc.edu.eg/apps/.../file.pdf"><br>
        <button type="button" onclick="testExtract()">Test Text Extraction</button>
      </form>
      <pre id="extractResult" style="border: 1px solid #ccc; padding: 10px; margin-top: 10px; white-space: pre-wrap; word-wrap: break-word;"></pre>

      <script>
      function testProxy() {
          var username = document.getElementById("username_proxy").value;
          var password = document.getElementById("password_proxy").value;
          var fileUrl = document.getElementById("fileUrl_proxy").value;
          var resultDiv = document.getElementById("proxyResult");
          resultDiv.innerText = "Fetching file...";

          var url = "/api/proxy?username=" + encodeURIComponent(username) +
                    "&password=" + encodeURIComponent(password) +
                    "&fileUrl=" + encodeURIComponent(fileUrl);
          fetch(url)
          .then(response => {
              if (!response.ok) {
                  return response.text().then(text => { throw new Error("Error " + response.status + ": " + (text || response.statusText)); });
              }
              // Get filename from URL or Content-Disposition if available (more complex)
              var filename = fileUrl.split('/').pop().split('?')[0] || 'downloaded_file';
              return response.blob().then(blob => ({ blob, filename }));
          })
          .then(({ blob, filename }) => {
              var downloadUrl = URL.createObjectURL(blob);
              resultDiv.innerHTML =
                "File ready. <a href='" + downloadUrl + "' download='" + filename + "'>Download '" + filename + "'</a>";
          })
          .catch(error => {
              resultDiv.innerText = "Error: " + error.message;
              console.error('Proxy fetch error:', error);
          });
      }

      function testExtract() {
          var username = document.getElementById("username_extract").value;
          var password = document.getElementById("password_extract").value;
          var fileUrl = document.getElementById("fileUrl_extract").value;
          var resultDiv = document.getElementById("extractResult");
          resultDiv.innerText = "Extracting text...";

          var url = "/api/extract?username=" + encodeURIComponent(username) +
                    "&password=" + encodeURIComponent(password) +
                    "&fileUrl=" + encodeURIComponent(fileUrl);
          fetch(url)
          .then(response => {
              if (!response.ok) {
                   return response.json().then(errData => { throw new Error("Error " + response.status + ": " + (errData.error || response.statusText)); });
              }
              return response.json();
          })
          .then(data => {
              resultDiv.innerText = data.text || "No text extracted.";
          })
          .catch(error => {
              resultDiv.innerText = "Error: " + error.message;
              console.error('Extract fetch error:', error);
          });
      }
      </script>
    </body>
    </html>
    """
    return render_template_string(html)


@app.route("/api/file-info", methods=["GET"])
def file_info():
    username = request.args.get("username")
    password = request.args.get("password")
    file_url = request.args.get("fileUrl")
    g.username = username # Set for logging

    if not username or not password or not file_url:
        g.log_outcome = "validation_error"
        g.log_error_message = "Missing required parameters"
        return jsonify({"error": "Missing required parameters"}), 400

    try:
        logger.info(f"Fetching HEAD for file info: {file_url}")
        response = session.head(
            file_url, auth=HttpNtlmAuth(username, password), timeout=15 # Increased timeout
        )

        if response.status_code != 200:
            g.log_outcome = "fetch_error_head"
            g.log_error_message = f"HEAD request failed: Status {response.status_code}"
            logger.warning(f"HEAD request failed for {file_url}: Status {response.status_code}")
            return jsonify({"error": f"Server returned {response.status_code}"}), response.status_code

        content_type = response.headers.get("Content-Type", "application/octet-stream")
        content_length = response.headers.get("Content-Length", "unknown")
        last_modified = response.headers.get("Last-Modified", "unknown")
        filename = file_url.split("/")[-1].split("?")[0] or "unknown" # Basic filename extraction

        g.log_outcome = "success"
        logger.info(f"Successfully retrieved file info for: {file_url}")
        return jsonify({
            "contentType": content_type,
            "contentLength": content_length,
            "lastModified": last_modified,
            "filename": filename,
        })

    except requests.exceptions.Timeout:
        g.log_outcome = "fetch_error_timeout"
        g.log_error_message = "Timeout during HEAD request"
        logger.error(f"Timeout during HEAD request for {file_url}")
        return jsonify({"error": "Timeout fetching file info"}), 504
    except requests.exceptions.RequestException as e:
        g.log_outcome = "fetch_error_network"
        g.log_error_message = f"Network error during HEAD request: {e}"
        logger.error(f"Network error during HEAD request for {file_url}: {e}")
        return jsonify({"error": f"Network error: {e}"}), 502 # Bad Gateway might be appropriate
    except Exception as e:
        g.log_outcome = "internal_error_unhandled"
        g.log_error_message = f"Unexpected error during file info: {e}"
        logger.exception(f"Unexpected error getting file info for {file_url}: {e}")
        return jsonify({"error": f"An unexpected error occurred: {e}"}), 500


@app.route("/api/proxy", methods=["GET"])
def proxy_file():
    bot_param = request.args.get("bot")
    if bot_param and bot_param.lower() == "true":
        logging.info("Received bot health check request for attendance API.")
        g.log_outcome = "bot_check_success" # Set outcome for logging
        # No username needed for bot check
        return jsonify({"status": "Success", "message": "Proxy API route is up!", "data": None}), 200
    username = request.args.get("username")
    password = request.args.get("password")
    file_url = request.args.get("fileUrl")
    prefetch = request.args.get("prefetch", "false").lower() == "true"
    g.username = username # Set for logging

    if not username or not password or not file_url:
        g.log_outcome = "validation_error"
        g.log_error_message = "Missing required parameters"
        logger.error("Proxy request missing required parameter(s).")
        # Use jsonify for consistent error responses
        return jsonify({"error": "Missing one or more required parameters: username, password, fileUrl"}), 400

    cache_key = get_cache_key(file_url, username)

    if prefetch:
        g.log_outcome = "prefetch_request"
        if ENABLE_BACKGROUND_THREADS and redis_client:
            with prefetch_lock:
                # Avoid adding duplicates
                if (username, password, file_url) not in prefetch_queue:
                     prefetch_queue.append((username, password, file_url))
                     logger.info(f"Prefetch queued for: {file_url}")
                     g.log_outcome = "prefetch_queued"
                else:
                     logger.info(f"Prefetch request already in queue for: {file_url}")
                     g.log_outcome = "prefetch_already_queued"
            return jsonify({"status": "prefetch queued"}), 200
        else:
            logger.info("Prefetch request received but background threads disabled or Redis unavailable.")
            g.log_outcome = "prefetch_unavailable"
            return jsonify({"status": "prefetch not available"}), 200 # Still success, just didn't queue

    range_header = request.headers.get("Range")

    # --- Cache Check ---
    # Skip cache check if range header is present, as we need fresh headers from server
    if not range_header and is_cached(cache_key):
        logger.info(f"Cache hit for: {file_url}")
        try:
            content = get_from_cache(cache_key)
            if content:
                file_name = file_url.split("/")[-1].split("?")[0] or "downloaded_file"
                content_type = guess_content_type(file_name)
                headers = {
                    "Content-Disposition": f'attachment; filename="{file_name}"',
                    "Content-Type": content_type,
                    "Content-Length": str(len(content)),
                    "Cache-Control": "public, max-age=86400", # Allow client caching
                    "X-Source": "redis-cache",
                    "Accept-Ranges": "bytes", # Indicate range support
                }
                g.log_outcome = "cache_hit"
                logger.info(f"Sending cached file from Redis ({len(content)} bytes) for {file_url}")
                return Response(content, headers=headers, status=200)
            else:
                 logger.warning(f"Cache hit reported for {cache_key}, but get_from_cache returned None. Cache might be inconsistent.")
                 # Proceed to download
        except Exception as e:
            logger.error(f"Error reading from Redis cache for {cache_key}: {e}", exc_info=True)
            # Proceed to download

    # --- Download (Cache Miss or Range Request) ---
    logger.info(f"Cache miss or range request, downloading: {file_url}")
    g.log_outcome = "fetch_attempt"
    try:
        req_headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
            "Accept": "*/*",
            "Accept-Encoding": "gzip, deflate", # Let requests handle decompression
        }
        if range_header:
            req_headers["Range"] = range_header
            logger.info(f"Requesting range: {range_header}")

        response = session.get(
            file_url,
            auth=HttpNtlmAuth(username, password),
            stream=True, # Crucial for handling large files and streaming
            timeout=60, # Generous timeout for downloads
            headers=req_headers,
        )

        # Check for auth errors specifically (often 401)
        if response.status_code == 401:
             g.log_outcome = "fetch_error_auth"
             g.log_error_message = "Authentication failed (401)"
             logger.warning(f"Authentication failed (401) for {file_url}")
             return jsonify({"error": "Authentication failed"}), 401

        # Check for other non-success codes
        if response.status_code not in (200, 206): # 206 Partial Content
            g.log_outcome = "fetch_error_server"
            g.log_error_message = f"Server error: Status {response.status_code}"
            logger.error(f"Server error {response.status_code} for {file_url}: {response.reason}")
            # Try to get error message from body if possible
            error_body = response.text[:500] if response.content else response.reason
            return jsonify({"error": f"Error from server ({response.status_code}): {error_body}"}), response.status_code

        # --- Successful Fetch ---
        content_type = response.headers.get("Content-Type", "application/octet-stream")
        file_name = file_url.split("/")[-1].split("?")[0] or "downloaded_file"
        content_length_str = response.headers.get("Content-Length") # Keep as string initially

        # --- Handle Range Request Response (Status 206) ---
        if response.status_code == 206:
            g.log_outcome = "fetch_partial_success"
            logger.info(f"Serving partial content (206) for {file_name}")
            def generate_partial():
                try:
                    for chunk in response.iter_content(chunk_size=CHUNK_SIZE):
                        yield chunk
                except Exception as e:
                    logger.error(f"Streaming error (partial content): {e}")
                    # Cannot set g here as request context might be gone
            resp_headers = {
                "Content-Type": content_type,
                "Content-Range": response.headers.get("Content-Range", ""),
                "Content-Length": content_length_str or "",
                "Accept-Ranges": "bytes",
            }
            return Response(generate_partial(), headers=resp_headers, status=206)

        # --- Handle Full Request Response (Status 200) ---
        g.log_outcome = "fetch_full_success"
        content_length = int(content_length_str) if content_length_str else 0

        # --- Streaming Strategy ---
        # Always stream directly for efficiency, cache in background if enabled/possible
        logger.info(f"Streaming full file ({content_length_str or 'unknown'} bytes) for {file_name}")
        def generate_full():
            try:
                # Use a temporary buffer to store content for potential caching
                # Only buffer if caching is intended and possible
                should_cache = redis_client and not is_cached(cache_key)
                buffered_content = bytearray() if should_cache else None

                for chunk in response.iter_content(chunk_size=CHUNK_SIZE):
                    yield chunk
                    if should_cache:
                        buffered_content.extend(chunk)

                # After streaming finishes, save to cache if buffered
                if should_cache and buffered_content:
                    logger.info(f"Finished streaming, attempting to cache {len(buffered_content)} bytes for {cache_key}")
                    # Run caching in a separate thread if enabled, otherwise block briefly
                    final_content = bytes(buffered_content)
                    if ENABLE_BACKGROUND_THREADS:
                        threading.Thread(target=save_to_cache, args=(cache_key, final_content), daemon=True).start()
                    else:
                        # In serverless, block briefly to cache (might timeout on large files)
                        save_to_cache(cache_key, final_content)

            except Exception as e:
                logger.error(f"Streaming error (full content): {e}")
                # Cannot set g here

        resp_headers = {
            "Content-Disposition": f'attachment; filename="{file_name}"',
            "Content-Type": content_type,
            "Accept-Ranges": "bytes", # Always accept ranges
            "Cache-Control": "public, max-age=86400", # Allow client caching
        }
        if content_length_str:
            resp_headers["Content-Length"] = content_length_str

        return Response(generate_full(), headers=resp_headers, status=200)

    except requests.exceptions.Timeout:
        g.log_outcome = "fetch_error_timeout"
        g.log_error_message = "Timeout during file download"
        logger.error(f"Request timed out downloading {file_url}")
        return jsonify({"error": "Timeout error while fetching file."}), 504
    except requests.exceptions.RequestException as req_err:
        g.log_outcome = "fetch_error_network"
        g.log_error_message = f"Network error during file download: {req_err}"
        logger.error(f"Request error downloading {file_url}: {req_err}")
        return jsonify({"error": f"Error fetching file: {req_err}"}), 502
    except Exception as e:
        g.log_outcome = "internal_error_unhandled"
        g.log_error_message = f"Unexpected error during proxy: {e}"
        logger.exception(f"Unexpected error in proxy for {file_url}: {e}")
        return jsonify({"error": f"An unexpected error occurred: {e}"}), 500


# Helper function to guess content type
def guess_content_type(filename):
    ext = filename.lower().split(".")[-1].split("?")[0] if "." in filename else ""
    content_types = {
        "pdf": "application/pdf", "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "doc": "application/msword", "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "ppt": "application/vnd.ms-powerpoint", "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "xls": "application/vnd.ms-excel", "txt": "text/plain", "jpg": "image/jpeg",
        "jpeg": "image/jpeg", "png": "image/png", "gif": "image/gif", "zip": "application/zip",
        "rar": "application/vnd.rar", "csv": "text/csv", "html": "text/html", "htm": "text/html",
    }
    return content_types.get(ext, "application/octet-stream")


@app.route("/api/extract", methods=["GET"])
def extract_text():
    username = request.args.get("username")
    password = request.args.get("password")
    file_url = request.args.get("fileUrl")
    g.username = username # Set for logging

    if not username or not password or not file_url:
        g.log_outcome = "validation_error"
        g.log_error_message = "Missing required parameters"
        logger.error("Extraction request missing required parameter(s).")
        return jsonify({"error": "Missing one or more required parameters: username, password, fileUrl"}), 400

    logger.info(f"Starting text extraction for: {file_url}")
    g.log_outcome = "extract_attempt"
    try:
        # Check cache first
        cache_key = get_cache_key(file_url, username)
        content = None
        if is_cached(cache_key):
             logger.info(f"Cache hit for extraction: {file_url}")
             content = get_from_cache(cache_key)
             if content:
                 g.log_outcome = "extract_cache_hit"
             else:
                  logger.warning(f"Cache hit reported for extraction {cache_key}, but get_from_cache failed. Fetching.")
                  g.log_outcome = "extract_cache_fail_fetch"

        if not content:
            logger.info(f"Cache miss or failed read for extraction, fetching: {file_url}")
            g.log_outcome = "extract_fetch_attempt"
            response = session.get(
                file_url, auth=HttpNtlmAuth(username, password), timeout=45 # Timeout for download
            )
            if response.status_code == 401:
                g.log_outcome = "extract_fetch_error_auth"
                g.log_error_message = "Authentication failed (401)"
                logger.warning(f"Authentication failed (401) for extraction: {file_url}")
                return jsonify({"error": "Authentication failed"}), 401
            if response.status_code != 200:
                g.log_outcome = "extract_fetch_error_server"
                g.log_error_message = f"Server error: Status {response.status_code}"
                logger.error(f"Server error {response.status_code} fetching for extraction: {response.text[:500]}")
                return jsonify({"error": f"Error from server ({response.status_code})"}), response.status_code

            content = response.content
            # Cache the downloaded content if Redis is available
            if redis_client:
                save_to_cache(cache_key, content)

        # --- Extraction Logic ---
        file_extension = file_url.split("/")[-1].split(".")[-1].split("?")[0].lower()
        logger.info(f"Extracting content ({len(content)} bytes) with extension: {file_extension}")
        extracted_text = ""
        extraction_method = "unsupported"

        if file_extension == "pdf":
            extraction_method = "pypdf2"
            pdf_bytes = BytesIO(content)
            try:
                reader = PyPDF2.PdfReader(pdf_bytes)
                num_pages = len(reader.pages)
                logger.info(f"PDF has {num_pages} pages.")
                page_texts = []
                for i in range(num_pages):
                    try:
                        page = reader.pages[i]
                        text = page.extract_text()
                        page_texts.append(text if text else "")
                    except Exception as page_err:
                         logger.warning(f"Error extracting text from PDF page {i+1}: {page_err}")
                         page_texts.append("") # Add empty string for failed pages
                extracted_text = "\n".join(page_texts)
                if not extracted_text.strip() and num_pages > 0:
                     logger.warning("PyPDF2 extracted no text, trying pdfminer fallback.")
                     # Fallback logic moved outside the main try block for clarity
            except Exception as e:
                logger.error(f"PyPDF2 failed for {file_url}: {e}. Will try pdfminer fallback.")
                # Fallback logic will trigger below if extracted_text is still empty

            # pdfminer fallback if PyPDF2 failed or yielded no text
            if not extracted_text.strip() and len(content) > 0:
                 try:
                     from pdfminer.high_level import extract_text as pdfminer_extract
                     logger.info("Using pdfminer.six fallback for PDF extraction.")
                     extraction_method = "pdfminer"
                     pdf_bytes.seek(0) # Reset stream position
                     extracted_text = pdfminer_extract(pdf_bytes)
                 except ImportError:
                      logger.error("pdfminer.six not installed, cannot use fallback.")
                      extracted_text = "PDF extraction failed (PyPDF2 error, pdfminer not installed)."
                 except Exception as pdfminer_err:
                      logger.error(f"pdfminer.six fallback failed: {pdfminer_err}")
                      extracted_text = f"PDF extraction failed (PyPDF2 and pdfminer errors: {pdfminer_err})."

        elif file_extension == "docx":
            extraction_method = "python-docx"
            try:
                import docx
                doc = docx.Document(BytesIO(content))
                paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
                extracted_text = "\n".join(paragraphs)
            except ImportError:
                 logger.error("python-docx not installed, cannot extract DOCX.")
                 extracted_text = "DOCX extraction failed (python-docx not installed)."
            except Exception as docx_err:
                 logger.error(f"python-docx failed: {docx_err}")
                 extracted_text = f"DOCX extraction failed: {docx_err}."

        elif file_extension == "pptx":
            extraction_method = "python-pptx"
            try:
                from pptx import Presentation
                prs = Presentation(BytesIO(content))
                slides_text = []
                for slide in prs.slides:
                    slide_text_parts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text and shape.text.strip():
                            slide_text_parts.append(shape.text.strip())
                    if slide_text_parts:
                        slides_text.append("\n".join(slide_text_parts))
                extracted_text = "\n\n".join(slides_text) # Separate slides
            except ImportError:
                 logger.error("python-pptx not installed, cannot extract PPTX.")
                 extracted_text = "PPTX extraction failed (python-pptx not installed)."
            except Exception as pptx_err:
                 logger.error(f"python-pptx failed: {pptx_err}")
                 extracted_text = f"PPTX extraction failed: {pptx_err}."
        else:
            logger.warning(f"Unsupported file type for extraction: {file_extension}")
            extracted_text = "Unsupported file type for extraction."
            g.log_outcome = "extract_unsupported_type"
            g.log_error_message = f"Unsupported file type: {file_extension}"
            # Return early for unsupported types? Or return the message? Let's return message.

        logger.info(f"Extraction complete using {extraction_method}. Extracted {len(extracted_text)} chars.")
        if not extracted_text.strip() and extraction_method != "unsupported":
             logger.warning(f"Extraction yielded no text for supported type {file_extension}.")
             g.log_outcome = "extract_success_no_text"
        elif extraction_method != "unsupported":
             g.log_outcome = "extract_success"

        return jsonify({"text": extracted_text.strip()}) # Return stripped text

    except requests.exceptions.Timeout:
        g.log_outcome = "extract_fetch_error_timeout"
        g.log_error_message = "Timeout during file download for extraction"
        logger.error(f"Timeout fetching for extraction: {file_url}")
        return jsonify({"error": "Timeout fetching file for extraction"}), 504
    except requests.exceptions.RequestException as req_err:
        g.log_outcome = "extract_fetch_error_network"
        g.log_error_message = f"Network error during extraction download: {req_err}"
        logger.error(f"Network error fetching for extraction: {file_url}: {req_err}")
        return jsonify({"error": f"Network error fetching file: {req_err}"}), 502
    except Exception as e:
        g.log_outcome = "extract_error_unhandled"
        g.log_error_message = f"Unexpected error during extraction: {e}"
        logger.exception(f"Unexpected error during extraction for {file_url}: {e}")
        return jsonify({"error": f"An unexpected error occurred during extraction: {e}"}), 500


# --- Main Execution ---
if __name__ == "__main__":
    # Add shutdown hook for the *logging* executor
    def shutdown_log_executor():
        print("Shutting down log executor...")
        log_executor.shutdown(wait=True)
        print("Log executor shut down.")
    atexit.register(shutdown_log_executor)

    # Use Waitress or Gunicorn for production
    logger.info(f"Starting Flask app for Proxy/Extractor in {'DEBUG' if app.debug else 'PRODUCTION'} mode.")
    # Example: from waitress import serve
    #          serve(app, host='0.0.0.0', port=5000, threads=8)
    app.run(host="0.0.0.0", port=5000, debug=app.debug)